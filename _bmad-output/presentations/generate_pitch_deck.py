"""
ReformLab Pitch Deck Generator
Generates a PPTX following the ReformLab presentation design prompt exactly.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import os

# ==============================================================================
# COLOR PALETTE (from design prompt)
# ==============================================================================
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE_50 = RGBColor(0xF8, 0xFA, 0xFC)
SLATE_200 = RGBColor(0xE2, 0xE8, 0xF0)
SLATE_300 = RGBColor(0xCB, 0xD5, 0xE1)
SLATE_400 = RGBColor(0x94, 0xA3, 0xB8)
SLATE_500 = RGBColor(0x64, 0x74, 0x8B)
SLATE_700 = RGBColor(0x33, 0x41, 0x55)
SLATE_900 = RGBColor(0x0F, 0x17, 0x2A)
BLUE_50 = RGBColor(0xEF, 0xF6, 0xFF)
BLUE_500 = RGBColor(0x3B, 0x82, 0xF6)
VIOLET_500 = RGBColor(0x8B, 0x5C, 0xF6)
EMERALD_500 = RGBColor(0x10, 0xB9, 0x81)
AMBER_500 = RGBColor(0xF5, 0x9E, 0x0B)
RED_500 = RGBColor(0xEF, 0x44, 0x44)

# Fonts (fallback-friendly)
FONT_HEADING = "Calibri"  # Inter fallback
FONT_BODY = "Calibri"
FONT_MONO = "Consolas"  # IBM Plex Mono fallback

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Margins
MARGIN = Inches(0.6)
CONTENT_LEFT = MARGIN
CONTENT_TOP = Inches(1.5)
CONTENT_WIDTH = SLIDE_WIDTH - 2 * MARGIN
CONTENT_HEIGHT = SLIDE_HEIGHT - CONTENT_TOP - MARGIN


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank

    # =========================================================================
    # SLIDE 1: TITLE
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    # Accent line at top
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), BLUE_500)

    # Product name
    add_text_box(
        slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
        "ReformLab", FONT_HEADING, Pt(48), SLATE_900, bold=True,
        alignment=PP_ALIGN.CENTER
    )

    # Tagline
    add_text_box(
        slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.8),
        "See the impact before the vote.", FONT_HEADING, Pt(28), BLUE_500,
        alignment=PP_ALIGN.CENTER, italic=True
    )

    # Subtitle
    add_text_box(
        slide, Inches(2.5), Inches(4.3), Inches(8), Inches(0.8),
        "Simulate environmental reforms. Understand who wins, who loses.\nTrust the results.",
        FONT_BODY, Pt(18), SLATE_500, alignment=PP_ALIGN.CENTER
    )

    # Visual element: stylized before/after bar chart showing households
    chart_top = Inches(5.3)
    chart_h = Inches(1.2)
    bar_w = Inches(0.5)
    gap = Inches(0.15)
    start_x = Inches(4.0)

    # "Before" bars (neutral)
    heights_before = [0.8, 0.6, 0.9, 0.5, 0.7]
    for i, h in enumerate(heights_before):
        x = start_x + i * (bar_w + gap)
        y = chart_top + chart_h - Inches(h)
        add_rect(slide, x, y, bar_w, Inches(h), SLATE_300, border=False)

    # Arrow
    arrow_x = start_x + 5 * (bar_w + gap) + Inches(0.3)
    add_text_box(slide, arrow_x, chart_top + Inches(0.35), Inches(0.8), Inches(0.5),
                 "\u2192", FONT_HEADING, Pt(32), SLATE_400, alignment=PP_ALIGN.CENTER)

    # "After" bars (colored by impact)
    after_colors = [EMERALD_500, BLUE_500, EMERALD_500, AMBER_500, RED_500]
    heights_after = [1.0, 0.7, 0.85, 0.6, 0.4]
    start_after = arrow_x + Inches(1.1)
    for i, (h, c) in enumerate(zip(heights_after, after_colors)):
        x = start_after + i * (bar_w + gap)
        y = chart_top + chart_h - Inches(h)
        add_rect(slide, x, y, bar_w, Inches(h), c, border=False)

    # Labels under bar groups
    add_text_box(slide, start_x, chart_top + chart_h + Inches(0.05),
                 Inches(3.2), Inches(0.4), "Current policy",
                 FONT_BODY, Pt(14), SLATE_400, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, start_after, chart_top + chart_h + Inches(0.05),
                 Inches(3.2), Inches(0.4), "After reform",
                 FONT_BODY, Pt(14), SLATE_400, alignment=PP_ALIGN.CENTER)

    add_speaker_notes(slide,
        "Key talking point: ReformLab lets you see the distributional impact of environmental reforms before they're enacted.\n\n"
        "Emotional beat: Curiosity — 'What is this and why should I care?'\n\n"
        "Transition: 'Let me show you what environmental reforms actually look like today.'")

    # =========================================================================
    # SLIDE 2: THE CONTEXT
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_slide_title(slide, "Governments are transforming how we live, move, and consume energy")

    # Four domain cards in a 2x2 grid
    domains = [
        ("\u26A1", "Energy Pricing", "Carbon taxes, energy tariffs,\nfuel levies", BLUE_500),
        ("\U0001F697", "Transport", "Emission standards, low-emission\nzones, mobility subsidies", VIOLET_500),
        ("\U0001F3E0", "Housing", "Renovation incentives, heating\nregulations, efficiency mandates", EMERALD_500),
        ("\u2696\uFE0F", "Redistribution", "Green dividends, targeted\nrebates, compensation schemes", AMBER_500),
    ]

    card_w = Inches(5.5)
    card_h = Inches(1.8)
    col_gap = Inches(0.6)
    row_gap = Inches(0.4)
    grid_left = (SLIDE_WIDTH - 2 * card_w - col_gap) / 2

    for i, (icon, title, desc, color) in enumerate(domains):
        col = i % 2
        row = i // 2
        x = grid_left + col * (card_w + col_gap)
        y = CONTENT_TOP + Inches(0.2) + row * (card_h + row_gap)

        # Card background
        add_rounded_rect(slide, x, y, card_w, card_h, SLATE_50, border_color=SLATE_200)

        # Color accent bar on left
        add_rect(slide, x, y, Inches(0.08), card_h, color, border=False)

        # Icon
        add_text_box(slide, x + Inches(0.3), y + Inches(0.25), Inches(0.7), Inches(0.7),
                     icon, FONT_HEADING, Pt(28), color, alignment=PP_ALIGN.CENTER)

        # Title
        add_text_box(slide, x + Inches(1.1), y + Inches(0.2), Inches(4.0), Inches(0.5),
                     title, FONT_HEADING, Pt(20), SLATE_900, bold=True)

        # Description
        add_text_box(slide, x + Inches(1.1), y + Inches(0.75), Inches(4.0), Inches(0.8),
                     desc, FONT_BODY, Pt(16), SLATE_500)

    # Bottom question
    add_text_box(
        slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.6),
        "Who pays, who benefits, and by how much?",
        FONT_HEADING, Pt(22), SLATE_900, bold=True, alignment=PP_ALIGN.CENTER
    )

    add_speaker_notes(slide,
        "Key talking point: Environmental reforms aren't abstract — they touch every household's budget through energy costs, transport, housing, and redistribution.\n\n"
        "Emotional beat: 'Oh, these are real policies affecting real people.'\n\n"
        "Transition: 'So who actually needs answers to these questions?'")

    # =========================================================================
    # SLIDE 3: THE AUDIENCE
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_slide_title(slide, "Three groups need rigorous, fast environmental policy assessment")

    personas = [
        {
            "icon": "\U0001F4CA",
            "title": "Policy Analyst",
            "context": "Government, consulting, think tank",
            "need": "Assess reforms before enactment",
            "pain": "Weeks of manual work per assessment",
            "want": "Fast, reproducible, briefing-ready",
            "color": BLUE_500,
        },
        {
            "icon": "\U0001F52C",
            "title": "Researcher",
            "context": "University or research institute",
            "need": "Project effects on real populations",
            "pain": "Custom throwaway code for each paper",
            "want": "Clean, reproducible projections",
            "color": VIOLET_500,
        },
        {
            "icon": "\U0001F465",
            "title": "Engaged Citizen",
            "context": "Voter, activist, journalist",
            "need": "Understand how reforms affect them",
            "pain": "Abstract numbers, no personal relevance",
            "want": "\"Show me what this means for me\"",
            "color": EMERALD_500,
        },
    ]

    card_w = Inches(3.7)
    card_h = Inches(4.5)
    total_cards_w = 3 * card_w + 2 * Inches(0.4)
    start_x = (SLIDE_WIDTH - total_cards_w) / 2

    for i, p in enumerate(personas):
        x = start_x + i * (card_w + Inches(0.4))
        y = CONTENT_TOP + Inches(0.2)

        # Card
        add_rounded_rect(slide, x, y, card_w, card_h, SLATE_50, border_color=SLATE_200)

        # Color top bar
        add_rect(slide, x, y, card_w, Inches(0.06), p["color"], border=False)

        # Icon
        add_text_box(slide, x, y + Inches(0.2), card_w, Inches(0.6),
                     p["icon"], FONT_HEADING, Pt(32), p["color"], alignment=PP_ALIGN.CENTER)

        # Title
        add_text_box(slide, x + Inches(0.3), y + Inches(0.85), card_w - Inches(0.6), Inches(0.4),
                     p["title"], FONT_HEADING, Pt(20), SLATE_900, bold=True, alignment=PP_ALIGN.CENTER)

        # Context
        add_text_box(slide, x + Inches(0.3), y + Inches(1.3), card_w - Inches(0.6), Inches(0.4),
                     p["context"], FONT_BODY, Pt(14), SLATE_500, alignment=PP_ALIGN.CENTER)

        # Divider
        add_rect(slide, x + Inches(0.5), y + Inches(1.8), card_w - Inches(1.0), Inches(0.02), SLATE_200, border=False)

        # Need
        add_text_box(slide, x + Inches(0.3), y + Inches(1.95), card_w - Inches(0.6), Inches(0.45),
                     p["need"], FONT_BODY, Pt(16), SLATE_700, alignment=PP_ALIGN.CENTER)

        # Pain (in muted red area)
        add_rounded_rect(slide, x + Inches(0.2), y + Inches(2.55), card_w - Inches(0.4), Inches(0.7),
                         RGBColor(0xFE, 0xF2, 0xF2), border_color=RGBColor(0xFE, 0xCA, 0xCA))
        add_text_box(slide, x + Inches(0.35), y + Inches(2.6), card_w - Inches(0.7), Inches(0.6),
                     p["pain"], FONT_BODY, Pt(14), RED_500, alignment=PP_ALIGN.CENTER)

        # Want (in subtle green area)
        add_rounded_rect(slide, x + Inches(0.2), y + Inches(3.45), card_w - Inches(0.4), Inches(0.7),
                         RGBColor(0xEC, 0xFD, 0xF5), border_color=RGBColor(0xA7, 0xF3, 0xD0))
        add_text_box(slide, x + Inches(0.35), y + Inches(3.5), card_w - Inches(0.7), Inches(0.6),
                     p["want"], FONT_BODY, Pt(14), EMERALD_500, alignment=PP_ALIGN.CENTER)

    add_speaker_notes(slide,
        "Key talking point: We serve three distinct audiences with different workflows but the same fundamental need — understanding who is affected by environmental reforms.\n\n"
        "Emotional beat: 'That's me / that's my team.'\n\n"
        "Transition: 'Let's look at what these people actually face today.'")

    # =========================================================================
    # SLIDE 4: THE PROBLEM
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_slide_title(slide, "Every environmental policy assessment starts from scratch")

    steps = [
        ("\u23F0", "Find and prepare data", "Scattered sources, incompatible formats"),
        ("\U0001F527", "Build the simulation", "Custom scripts rebuilt every time"),
        ("\U0001F504", "Run projections", "Manual year-by-year looping"),
        ("\U0001F4C9", "Compute winners and losers", "Custom analysis code each time"),
        ("\U0001F4DD", "Document methodology", "Usually skipped entirely"),
    ]

    step_h = Inches(0.8)
    step_gap = Inches(0.15)
    steps_left = Inches(1.5)
    steps_top = CONTENT_TOP + Inches(0.2)

    for i, (icon, title, desc) in enumerate(steps):
        y = steps_top + i * (step_h + step_gap)

        # Step number circle
        circle_size = Inches(0.55)
        add_circle(slide, steps_left, y + Inches(0.12), circle_size, SLATE_50, SLATE_300)
        add_text_box(slide, steps_left, y + Inches(0.12), circle_size, circle_size,
                     str(i + 1), FONT_MONO, Pt(18), SLATE_700, alignment=PP_ALIGN.CENTER,
                     vertical=MSO_ANCHOR.MIDDLE)

        # Connector line (except last)
        if i < len(steps) - 1:
            line_x = steps_left + circle_size / 2
            line_y = y + Inches(0.12) + circle_size
            add_rect(slide, line_x - Inches(0.015), line_y, Inches(0.03),
                     step_gap + Inches(0.0), SLATE_300, border=False)

        # Icon
        add_text_box(slide, steps_left + Inches(0.8), y + Inches(0.05), Inches(0.5), Inches(0.5),
                     icon, FONT_HEADING, Pt(22), SLATE_500)

        # Title
        add_text_box(slide, steps_left + Inches(1.5), y + Inches(0.05), Inches(4.5), Inches(0.4),
                     title, FONT_HEADING, Pt(18), SLATE_900, bold=True)

        # Description
        add_text_box(slide, steps_left + Inches(1.5), y + Inches(0.42), Inches(4.5), Inches(0.35),
                     desc, FONT_BODY, Pt(15), SLATE_500)

    # Result box on the right side
    result_x = Inches(7.5)
    result_y = CONTENT_TOP + Inches(0.3)
    result_w = Inches(4.8)
    result_h = Inches(3.5)

    add_rounded_rect(slide, result_x, result_y, result_w, result_h,
                     RGBColor(0xFE, 0xF2, 0xF2), border_color=RGBColor(0xFE, 0xCA, 0xCA))

    add_text_box(slide, result_x + Inches(0.4), result_y + Inches(0.3), result_w - Inches(0.8), Inches(0.5),
                 "Result:", FONT_HEADING, Pt(18), RED_500, bold=True)

    result_items = [
        "Weeks of work",
        "Fragile pipelines",
        "Undocumented assumptions",
        "No reproducibility",
    ]
    for j, item in enumerate(result_items):
        add_text_box(slide, result_x + Inches(0.4), result_y + Inches(0.9) + j * Inches(0.5),
                     result_w - Inches(0.8), Inches(0.45),
                     "\u2717  " + item, FONT_BODY, Pt(16), SLATE_700)

    # Gut-punch line at bottom
    add_text_box(
        slide, Inches(1.0), Inches(6.1), Inches(11), Inches(0.7),
        "\"What about low-income rural households?\" \u2014 You start over.",
        FONT_HEADING, Pt(20), RED_500, bold=True, alignment=PP_ALIGN.CENTER,
        italic=True
    )

    add_speaker_notes(slide,
        "Key talking point: Walk through each step — data, simulation, projection, analysis, documentation. Every single one is manual, fragile, and rebuilt from scratch.\n\n"
        "Emotional beat: Pain — 'Yes, this is exactly what we deal with.'\n\n"
        "Transition: 'You might ask — isn't there already a tool that does this?'")

    # =========================================================================
    # SLIDE 5: THE GAP
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_slide_title(slide, "No existing tool covers the full workflow")

    # Gap analysis table
    headers = ["Capability", "Tax-benefit\ntools", "Energy/building\ntools", "Ad-hoc\nscripts"]
    rows = [
        ["Cross-domain (income +\nenergy + transport)", "No", "No", "Fragile"],
        ["Multi-year projections", "No", "No", "Manual"],
        ["Distributional impact\n(who wins / loses)", "Partial", "No", "Manual"],
        ["Built-in reproducibility", "No", "No", "No"],
        ["No-code access", "No", "No", "No"],
        ["Environmental policy focus", "No", "No", "Rebuilt\neach time"],
    ]

    table_left = Inches(1.2)
    table_top = CONTENT_TOP + Inches(0.1)
    col_widths = [Inches(3.5), Inches(2.5), Inches(2.5), Inches(2.5)]
    table_width = sum(w for w in col_widths)
    row_h = Inches(0.65)
    header_h = Inches(0.7)

    # Header row
    for j, (header, cw) in enumerate(zip(headers, col_widths)):
        x = table_left + sum(col_widths[k] for k in range(j))
        add_rounded_rect(slide, x, table_top, cw, header_h, SLATE_50 if j == 0 else SLATE_50,
                         border_color=SLATE_200)
        add_text_box(slide, x + Inches(0.15), table_top + Inches(0.05), cw - Inches(0.3), header_h - Inches(0.1),
                     header, FONT_HEADING, Pt(15), SLATE_700, bold=True,
                     alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT,
                     vertical=MSO_ANCHOR.MIDDLE)

    # Data rows
    for i, row in enumerate(rows):
        y = table_top + header_h + i * row_h
        for j, (cell, cw) in enumerate(zip(row, col_widths)):
            x = table_left + sum(col_widths[k] for k in range(j))
            # Color-code cells
            if j > 0:
                if cell == "No":
                    bg = RGBColor(0xFE, 0xF2, 0xF2)
                    fg = RED_500
                elif cell in ("Fragile", "Manual", "Partial", "Rebuilt\neach time"):
                    bg = RGBColor(0xFF, 0xFB, 0xEB)
                    fg = AMBER_500
                else:
                    bg = WHITE
                    fg = SLATE_700
            else:
                bg = WHITE
                fg = SLATE_900

            add_rect(slide, x, y, cw, row_h, bg, border_color=SLATE_200)
            font_name = FONT_MONO if j > 0 else FONT_BODY
            add_text_box(slide, x + Inches(0.15), y + Inches(0.02), cw - Inches(0.3), row_h - Inches(0.04),
                         cell, font_name, Pt(14), fg,
                         alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT,
                         vertical=MSO_ANCHOR.MIDDLE)

    # Conclusion
    add_text_box(
        slide, Inches(1.5), Inches(6.0), Inches(10), Inches(0.6),
        "There is no integrated platform for environmental reform assessment.",
        FONT_HEADING, Pt(20), SLATE_900, bold=True, alignment=PP_ALIGN.CENTER
    )

    add_speaker_notes(slide,
        "Key talking point: Walk through the table — tax-benefit tools handle income rules but not environment; energy tools handle physics but not distribution; scripts are fragile. Nobody covers all six rows.\n\n"
        "Emotional beat: Frustration — 'Right, nothing solves this properly.'\n\n"
        "Transition: 'That's exactly what ReformLab is built to solve.'")

    # =========================================================================
    # SLIDE 6: THE SOLUTION
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_slide_title(slide, "End-to-end environmental policy assessment \u2014 no code required")

    # Four-step horizontal flow
    flow_steps = [
        ("\u2699\uFE0F", "Configure", "Pick a template. Adjust\nparameters. Set timeline.", BLUE_500),
        ("\u25B6\uFE0F", "Simulate", "Run against a population.\nMulti-year projections.", VIOLET_500),
        ("\U0001F4CA", "Analyze", "Distributional impact.\nWinners and losers.", EMERALD_500),
        ("\u2705", "Trust", "Every assumption logged.\nFully reproducible.", AMBER_500),
    ]

    step_w = Inches(2.5)
    step_h_val = Inches(2.8)
    flow_gap = Inches(0.35)
    total_flow_w = 4 * step_w + 3 * flow_gap
    flow_left = (SLIDE_WIDTH - total_flow_w) / 2
    flow_top = CONTENT_TOP + Inches(0.3)

    for i, (icon, title, desc, color) in enumerate(flow_steps):
        x = flow_left + i * (step_w + flow_gap)

        # Card
        add_rounded_rect(slide, x, flow_top, step_w, step_h_val, WHITE, border_color=SLATE_200)

        # Color top bar
        add_rect(slide, x, flow_top, step_w, Inches(0.06), color, border=False)

        # Step number
        add_text_box(slide, x + Inches(0.15), flow_top + Inches(0.2), Inches(0.5), Inches(0.4),
                     str(i + 1), FONT_MONO, Pt(16), SLATE_400)

        # Icon
        add_text_box(slide, x, flow_top + Inches(0.55), step_w, Inches(0.6),
                     icon, FONT_HEADING, Pt(32), color, alignment=PP_ALIGN.CENTER)

        # Title
        add_text_box(slide, x + Inches(0.2), flow_top + Inches(1.2), step_w - Inches(0.4), Inches(0.45),
                     title, FONT_HEADING, Pt(20), SLATE_900, bold=True, alignment=PP_ALIGN.CENTER)

        # Description
        add_text_box(slide, x + Inches(0.2), flow_top + Inches(1.75), step_w - Inches(0.4), Inches(0.8),
                     desc, FONT_BODY, Pt(15), SLATE_500, alignment=PP_ALIGN.CENTER)

        # Arrow connector (except last)
        if i < len(flow_steps) - 1:
            arrow_x_pos = x + step_w + Inches(0.05)
            add_text_box(slide, arrow_x_pos, flow_top + Inches(1.1), Inches(0.25), Inches(0.5),
                         "\u2192", FONT_HEADING, Pt(24), SLATE_300, alignment=PP_ALIGN.CENTER)

    # Subtitle below flow
    add_text_box(
        slide, Inches(2.0), flow_top + step_h_val + Inches(0.4), Inches(9), Inches(0.6),
        "Configure a reform, run it against a population, see who wins and who loses.",
        FONT_BODY, Pt(18), SLATE_500, alignment=PP_ALIGN.CENTER
    )

    # GUI mockup hint bar
    mock_y = flow_top + step_h_val + Inches(1.2)
    add_rounded_rect(slide, Inches(2.5), mock_y, Inches(8), Inches(1.2),
                     SLATE_50, border_color=SLATE_200)

    # Mock header bar
    add_rect(slide, Inches(2.5), mock_y, Inches(8), Inches(0.35), SLATE_200, border=False)
    add_text_box(slide, Inches(2.7), mock_y + Inches(0.02), Inches(3), Inches(0.3),
                 "\u25CF \u25CF \u25CF   ReformLab \u2014 Scenario Comparison",
                 FONT_MONO, Pt(14), SLATE_500)

    # Mini chart placeholders inside mock
    mock_content_y = mock_y + Inches(0.4)
    bar_colors = [BLUE_500, VIOLET_500, EMERALD_500]
    for k, bc in enumerate(bar_colors):
        bx = Inches(3.0) + k * Inches(2.2)
        add_rect(slide, bx, mock_content_y + Inches(0.4), Inches(0.4), Inches(0.35), bc, border=False)
        add_rect(slide, bx + Inches(0.5), mock_content_y + Inches(0.25), Inches(0.4), Inches(0.5), bc, border=False)
        add_rect(slide, bx + Inches(1.0), mock_content_y + Inches(0.15), Inches(0.4), Inches(0.6), bc, border=False)

    add_speaker_notes(slide,
        "Key talking point: Walk through the four steps — configure, simulate, analyze, trust. Emphasize that the whole workflow is visual, no coding required.\n\n"
        "Emotional beat: Relief — 'Wait, this does everything in one place?'\n\n"
        "Transition: 'Let me show you specifically what makes this different from anything else out there.'")

    # =========================================================================
    # SLIDE 7: DIFFERENTIATORS
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_slide_title(slide, "What makes ReformLab different")

    diffs = [
        ("\U0001F5A5\uFE0F", "No-code, end-to-end workflow",
         "From policy definition to distributional charts\nwithout writing a single line of code.", BLUE_500),
        ("\U0001F4C4", "Environmental policy templates",
         "Pre-built templates for carbon taxes, subsidies,\ntransport policies. Configure, don't rebuild.", VIOLET_500),
        ("\U0001F4C8", "Dynamic multi-year projection",
         "See impacts evolve over 10+ years as fleets\nturn over and policies phase in.", EMERALD_500),
        ("\U0001F50D", "Assumption transparency",
         "Every parameter and data source is logged.\nFull reproducibility without extra effort.", AMBER_500),
        ("\U0001F310", "Open-data default",
         "Works out of the box with public data.\nBring your own when you have it.", SLATE_500),
    ]

    tile_h = Inches(0.85)
    tile_gap = Inches(0.2)
    tiles_left = Inches(1.0)
    tiles_top = CONTENT_TOP + Inches(0.2)
    tile_w = Inches(11.0)

    for i, (icon, title, desc, color) in enumerate(diffs):
        y = tiles_top + i * (tile_h + tile_gap)

        # Background
        bg_color = BLUE_50 if i == 0 else WHITE  # Highlight first one
        add_rounded_rect(slide, tiles_left, y, tile_w, tile_h, bg_color, border_color=SLATE_200)

        # Left color accent
        add_rect(slide, tiles_left, y, Inches(0.06), tile_h, color, border=False)

        # Icon
        add_text_box(slide, tiles_left + Inches(0.25), y + Inches(0.1), Inches(0.6), Inches(0.6),
                     icon, FONT_HEADING, Pt(24), color)

        # Title
        add_text_box(slide, tiles_left + Inches(1.0), y + Inches(0.08), Inches(4.5), Inches(0.4),
                     title, FONT_HEADING, Pt(18), SLATE_900, bold=True)

        # Description
        add_text_box(slide, tiles_left + Inches(5.8), y + Inches(0.08), Inches(5.0), Inches(0.7),
                     desc, FONT_BODY, Pt(15), SLATE_500)

    add_speaker_notes(slide,
        "Key talking point: Five differentiators — lead with no-code because that's the headline for this audience. Templates save months of work. Multi-year projection is unique. Transparency is built-in, not an afterthought.\n\n"
        "Emotional beat: Conviction — 'Okay, this is genuinely different.'\n\n"
        "Transition: 'Let me show you what this looks like for each of our three audiences.'")

    # =========================================================================
    # SLIDE 8: FOR THE POLICY ANALYST
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    # Quote at top
    add_rounded_rect(slide, MARGIN, Inches(0.4), CONTENT_WIDTH, Inches(0.9),
                     SLATE_50, border_color=SLATE_200)
    add_rect(slide, MARGIN, Inches(0.4), Inches(0.06), Inches(0.9), BLUE_500, border=False)
    add_text_box(slide, MARGIN + Inches(0.4), Inches(0.5), CONTENT_WIDTH - Inches(0.8), Inches(0.7),
                 "\u201CI modified one parameter and got updated charts instantly.\nThis used to take me weeks.\u201D",
                 FONT_HEADING, Pt(20), SLATE_700, italic=True, alignment=PP_ALIGN.CENTER)

    # Before/After split
    split_top = Inches(1.6)
    split_h = Inches(5.2)
    half_w = (CONTENT_WIDTH - Inches(0.3)) / 2

    # BEFORE side (left)
    add_rounded_rect(slide, MARGIN, split_top, half_w, split_h, SLATE_50, border_color=SLATE_200)
    add_text_box(slide, MARGIN + Inches(0.3), split_top + Inches(0.2), Inches(2), Inches(0.4),
                 "BEFORE", FONT_HEADING, Pt(16), SLATE_400, bold=True)

    before_items = [
        "\u2717  Chain three tools per assessment",
        "\u2717  Rebuild pipelines each quarter",
        "\u2717  Manual methodology documentation",
        "\u2717  Hours to answer follow-ups",
    ]
    for j, item in enumerate(before_items):
        add_text_box(slide, MARGIN + Inches(0.4), split_top + Inches(0.8) + j * Inches(0.55),
                     half_w - Inches(0.8), Inches(0.5),
                     item, FONT_BODY, Pt(16), SLATE_500)

    # Visual: cluttered workflow diagram (before)
    clutter_y = split_top + Inches(3.2)
    clutter_x = MARGIN + Inches(0.5)
    for k in range(3):
        bx = clutter_x + k * Inches(1.8)
        add_rounded_rect(slide, bx, clutter_y, Inches(1.4), Inches(0.6), WHITE, border_color=SLATE_300)
        add_text_box(slide, bx + Inches(0.1), clutter_y + Inches(0.1), Inches(1.2), Inches(0.4),
                     ["Tool A", "Tool B", "Tool C"][k], FONT_MONO, Pt(14), SLATE_400, alignment=PP_ALIGN.CENTER)
        if k < 2:
            add_text_box(slide, bx + Inches(1.45), clutter_y + Inches(0.1), Inches(0.3), Inches(0.4),
                         "\u2192", FONT_HEADING, Pt(18), RED_500)
    add_text_box(slide, clutter_x, clutter_y + Inches(0.7), Inches(5), Inches(0.4),
                 "\u26A0\uFE0F  Fragile chain, breaks often", FONT_BODY, Pt(14), RED_500)

    # AFTER side (right)
    after_x = MARGIN + half_w + Inches(0.3)
    add_rounded_rect(slide, after_x, split_top, half_w, split_h, WHITE, border_color=BLUE_500)
    add_text_box(slide, after_x + Inches(0.3), split_top + Inches(0.2), Inches(3), Inches(0.4),
                 "WITH REFORMLAB", FONT_HEADING, Pt(16), BLUE_500, bold=True)

    after_items = [
        "\u2713  Configure reform in the GUI",
        "\u2713  Run against built-in population",
        "\u2713  Distributional charts instantly",
        "\u2713  Re-run with alternatives in minutes",
    ]
    for j, item in enumerate(after_items):
        add_text_box(slide, after_x + Inches(0.4), split_top + Inches(0.8) + j * Inches(0.55),
                     half_w - Inches(0.8), Inches(0.5),
                     item, FONT_BODY, Pt(16), SLATE_900)

    # Visual: clean single-tool workflow (after)
    clean_y = split_top + Inches(3.2)
    clean_x = after_x + Inches(0.8)
    add_rounded_rect(slide, clean_x, clean_y, Inches(4.0), Inches(1.4), SLATE_50, border_color=BLUE_500)
    add_text_box(slide, clean_x, clean_y + Inches(0.05), Inches(4.0), Inches(0.35),
                 "ReformLab", FONT_HEADING, Pt(16), BLUE_500, bold=True, alignment=PP_ALIGN.CENTER)
    # Mini bars inside
    mini_colors = [BLUE_500, VIOLET_500, EMERALD_500, AMBER_500]
    for k, mc in enumerate(mini_colors):
        bx = clean_x + Inches(0.4) + k * Inches(0.9)
        bar_ht = [0.5, 0.7, 0.4, 0.6][k]
        add_rect(slide, bx, clean_y + Inches(1.0) - Inches(bar_ht * 0.6),
                 Inches(0.5), Inches(bar_ht * 0.6), mc, border=False)

    add_speaker_notes(slide,
        "Key talking point: The policy analyst's before/after story. Today it's a fragile chain of tools rebuilt every quarter. With ReformLab, configure once, run instantly, and answer follow-up questions in minutes, not hours.\n\n"
        "Emotional beat: Recognition — 'I see myself using this.'\n\n"
        "Transition: 'Now let's see what this means for researchers.'")

    # =========================================================================
    # SLIDE 9: FOR THE RESEARCHER
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    # Quote at top
    add_rounded_rect(slide, MARGIN, Inches(0.4), CONTENT_WIDTH, Inches(0.9),
                     SLATE_50, border_color=SLATE_200)
    add_rect(slide, MARGIN, Inches(0.4), Inches(0.06), Inches(0.9), VIOLET_500, border=False)
    add_text_box(slide, MARGIN + Inches(0.4), Inches(0.5), CONTENT_WIDTH - Inches(0.8), Inches(0.7),
                 "\u201CI projected my results onto a reproducible simulation \u2014\nmy co-author replicated them on a different machine.\u201D",
                 FONT_HEADING, Pt(20), SLATE_700, italic=True, alignment=PP_ALIGN.CENTER)

    # Before/After split
    split_top = Inches(1.6)
    split_h = Inches(5.2)
    half_w = (CONTENT_WIDTH - Inches(0.3)) / 2

    # BEFORE side
    add_rounded_rect(slide, MARGIN, split_top, half_w, split_h, SLATE_50, border_color=SLATE_200)
    add_text_box(slide, MARGIN + Inches(0.3), split_top + Inches(0.2), Inches(2), Inches(0.4),
                 "BEFORE", FONT_HEADING, Pt(16), SLATE_400, bold=True)

    before_items = [
        "\u2717  800 lines of custom code per paper",
        "\u2717  Results impossible to replicate",
        "\u2717  Co-authors can't run your simulation",
        "\u2717  Replication packages are incomplete",
    ]
    for j, item in enumerate(before_items):
        add_text_box(slide, MARGIN + Inches(0.4), split_top + Inches(0.8) + j * Inches(0.55),
                     half_w - Inches(0.8), Inches(0.5),
                     item, FONT_BODY, Pt(16), SLATE_500)

    # Code mess visual
    code_y = split_top + Inches(3.3)
    add_rounded_rect(slide, MARGIN + Inches(0.4), code_y, Inches(5.0), Inches(1.3),
                     RGBColor(0x1E, 0x29, 0x3B), border_color=SLATE_700)
    code_lines = [
        "# simulation_v3_final_FINAL.py",
        "for year in range(2025, 2035):",
        "    pop = load(f'data_{year}.csv')",
        "    # TODO: fix this hack",
    ]
    for k, line in enumerate(code_lines):
        add_text_box(slide, MARGIN + Inches(0.6), code_y + Inches(0.1) + k * Inches(0.28),
                     Inches(4.6), Inches(0.28),
                     line, FONT_MONO, Pt(14), SLATE_300)

    # AFTER side
    after_x = MARGIN + half_w + Inches(0.3)
    add_rounded_rect(slide, after_x, split_top, half_w, split_h, WHITE, border_color=VIOLET_500)
    add_text_box(slide, after_x + Inches(0.3), split_top + Inches(0.2), Inches(3), Inches(0.4),
                 "WITH REFORMLAB", FONT_HEADING, Pt(16), VIOLET_500, bold=True)

    after_items = [
        "\u2713  Python API + Jupyter notebooks",
        "\u2713  Multi-year projections built in",
        "\u2713  Immutable manifest per run",
        "\u2713  Reproducibility by default",
    ]
    for j, item in enumerate(after_items):
        add_text_box(slide, after_x + Inches(0.4), split_top + Inches(0.8) + j * Inches(0.55),
                     half_w - Inches(0.8), Inches(0.5),
                     item, FONT_BODY, Pt(16), SLATE_900)

    # Notebook mockup
    nb_y = split_top + Inches(3.3)
    add_rounded_rect(slide, after_x + Inches(0.4), nb_y, Inches(5.0), Inches(1.3),
                     SLATE_50, border_color=VIOLET_500)
    add_text_box(slide, after_x + Inches(0.6), nb_y + Inches(0.05), Inches(4), Inches(0.28),
                 "\U0001F4D3  ReformLab Notebook", FONT_MONO, Pt(14), VIOLET_500, bold=True)
    nb_lines = [
        "reform = CarbonTax(rate=44, start=2026)",
        "results = simulate(reform, years=10)",
        "results.distributional_impact()",
    ]
    for k, line in enumerate(nb_lines):
        add_text_box(slide, after_x + Inches(0.6), nb_y + Inches(0.38) + k * Inches(0.28),
                     Inches(4.6), Inches(0.28),
                     line, FONT_MONO, Pt(14), SLATE_700)

    add_speaker_notes(slide,
        "Key talking point: Researchers get full programmatic control via the Python API and Jupyter notebooks. Every run produces an immutable manifest. Reproducibility is built in, not bolted on.\n\n"
        "Emotional beat: Trust — 'This solves my reproducibility problem.'\n\n"
        "Transition: 'And for the broader public...'")

    # =========================================================================
    # SLIDE 10: FOR THE ENGAGED CITIZEN
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    # Quote at top
    add_rounded_rect(slide, MARGIN, Inches(0.4), CONTENT_WIDTH, Inches(0.7),
                     SLATE_50, border_color=SLATE_200)
    add_rect(slide, MARGIN, Inches(0.4), Inches(0.06), Inches(0.7), EMERALD_500, border=False)
    add_text_box(slide, MARGIN + Inches(0.4), Inches(0.45), CONTENT_WIDTH - Inches(0.8), Inches(0.5),
                 "\u201CI finally understand who wins and who loses.\u201D",
                 FONT_HEADING, Pt(22), SLATE_700, italic=True, alignment=PP_ALIGN.CENTER)

    # Citizen dashboard mockup (large)
    mock_top = Inches(1.4)
    mock_left = Inches(1.0)
    mock_w = Inches(11.0)
    mock_h = Inches(5.0)

    add_rounded_rect(slide, mock_left, mock_top, mock_w, mock_h, SLATE_50, border_color=SLATE_200)

    # Mock header
    add_rect(slide, mock_left, mock_top, mock_w, Inches(0.4), SLATE_200, border=False)
    add_text_box(slide, mock_left + Inches(0.3), mock_top + Inches(0.05), Inches(5), Inches(0.3),
                 "\u25CF \u25CF \u25CF   ReformLab \u2014 Compare Reform Proposals",
                 FONT_MONO, Pt(14), SLATE_500)

    # Two proposal columns inside mockup
    prop_top = mock_top + Inches(0.6)
    prop_w = Inches(5.0)
    prop_h = Inches(4.0)
    prop_gap = Inches(0.6)
    prop_left = mock_left + (mock_w - 2 * prop_w - prop_gap) / 2

    proposals = [
        ("Proposal A: Carbon Tax + Green Dividend", BLUE_500,
         [("Your household pays", "+\u20AC340/year", AMBER_500),
          ("Green dividend received", "\u2212\u20AC420/year", EMERALD_500),
          ("Net impact", "\u2212\u20AC80/year", EMERALD_500)]),
        ("Proposal B: Fuel Levy + Mobility Subsidy", VIOLET_500,
         [("Your household pays", "+\u20AC280/year", AMBER_500),
          ("Mobility subsidy received", "\u2212\u20AC150/year", EMERALD_500),
          ("Net impact", "+\u20AC130/year", RED_500)]),
    ]

    for i, (title, color, metrics) in enumerate(proposals):
        px = prop_left + i * (prop_w + prop_gap)

        add_rounded_rect(slide, px, prop_top, prop_w, prop_h, WHITE, border_color=color)
        add_rect(slide, px, prop_top, prop_w, Inches(0.06), color, border=False)

        # Proposal title
        add_text_box(slide, px + Inches(0.2), prop_top + Inches(0.2), prop_w - Inches(0.4), Inches(0.5),
                     title, FONT_HEADING, Pt(16), SLATE_900, bold=True)

        # Metrics
        for j, (label, value, val_color) in enumerate(metrics):
            my = prop_top + Inches(1.0) + j * Inches(0.8)
            add_text_box(slide, px + Inches(0.3), my, Inches(2.8), Inches(0.35),
                         label, FONT_BODY, Pt(15), SLATE_500)
            add_text_box(slide, px + Inches(3.0), my, Inches(1.7), Inches(0.35),
                         value, FONT_MONO, Pt(16), val_color, bold=True, alignment=PP_ALIGN.RIGHT)

            # Divider line after metric
            if j < len(metrics) - 1:
                add_rect(slide, px + Inches(0.3), my + Inches(0.45),
                         prop_w - Inches(0.6), Inches(0.015), SLATE_200, border=False)

        # Mini distributional chart at bottom of each proposal
        chart_y = prop_top + Inches(3.1)
        add_text_box(slide, px + Inches(0.3), chart_y - Inches(0.3), Inches(3), Inches(0.3),
                     "Impact by income quintile:", FONT_BODY, Pt(14), SLATE_500)
        quintile_heights = [[0.5, 0.35, 0.2, 0.15, 0.1], [0.3, 0.25, 0.2, 0.3, 0.4]]
        quintile_colors_a = [EMERALD_500, EMERALD_500, EMERALD_500, AMBER_500, AMBER_500]
        quintile_colors_b = [EMERALD_500, EMERALD_500, AMBER_500, RED_500, RED_500]
        q_colors = [quintile_colors_a, quintile_colors_b]

        for q in range(5):
            qx = px + Inches(0.4) + q * Inches(0.85)
            qh = quintile_heights[i][q]
            qy = chart_y + Inches(0.6) - Inches(qh)
            add_rect(slide, qx, qy, Inches(0.5), Inches(qh), q_colors[i][q], border=False)

        # Quintile labels
        add_text_box(slide, px + Inches(0.3), chart_y + Inches(0.65), prop_w - Inches(0.6), Inches(0.25),
                     "Q1        Q2        Q3        Q4        Q5",
                     FONT_MONO, Pt(14), SLATE_400, alignment=PP_ALIGN.CENTER)

    # Use case note at bottom
    add_text_box(
        slide, Inches(1.5), Inches(6.7), Inches(10), Inches(0.5),
        "Before elections: compare candidates' platforms. Before a vote: understand the real impact.",
        FONT_BODY, Pt(16), SLATE_500, alignment=PP_ALIGN.CENTER, italic=True
    )

    add_speaker_notes(slide,
        "Key talking point: A public web app built on ReformLab lets citizens enter their household profile and compare reform proposals side by side. They see personal impact and the big picture.\n\n"
        "Emotional beat: Inspiration — 'This could change how people understand policy.'\n\n"
        "Transition: 'You might be wondering — who else is doing this?'")

    # =========================================================================
    # SLIDE 11: COMPETITIVE LANDSCAPE
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_slide_title(slide, "The field is clearing \u2014 and no one occupies this space")

    # Positioning map (2x2 quadrant)
    map_left = Inches(1.0)
    map_top = CONTENT_TOP + Inches(0.3)
    map_w = Inches(7.5)
    map_h = Inches(5.0)

    # Quadrant background - highlight top-right
    # Top-right (ReformLab's quadrant)
    add_rect(slide, map_left + map_w / 2, map_top, map_w / 2, map_h / 2, BLUE_50, border=False)
    # Other quadrants
    add_rect(slide, map_left, map_top, map_w / 2, map_h / 2, SLATE_50, border=False)
    add_rect(slide, map_left, map_top + map_h / 2, map_w / 2, map_h / 2, SLATE_50, border=False)
    add_rect(slide, map_left + map_w / 2, map_top + map_h / 2, map_w / 2, map_h / 2, SLATE_50, border=False)

    # Axes
    # Horizontal axis
    add_rect(slide, map_left, map_top + map_h / 2 - Inches(0.015), map_w, Inches(0.03), SLATE_300, border=False)
    # Vertical axis
    add_rect(slide, map_left + map_w / 2 - Inches(0.015), map_top, Inches(0.03), map_h, SLATE_300, border=False)

    # Axis labels
    add_text_box(slide, map_left, map_top + map_h + Inches(0.1), map_w / 2, Inches(0.4),
                 "\u2190 Tax-benefit focused", FONT_BODY, Pt(14), SLATE_400, alignment=PP_ALIGN.LEFT)
    add_text_box(slide, map_left + map_w / 2, map_top + map_h + Inches(0.1), map_w / 2, Inches(0.4),
                 "Environmental focused \u2192", FONT_BODY, Pt(14), SLATE_400, alignment=PP_ALIGN.RIGHT)
    # Y-axis labels (rotated text is hard in pptx, use horizontal)
    add_text_box(slide, map_left - Inches(0.1), map_top - Inches(0.35), map_w / 2, Inches(0.3),
                 "Dynamic / No-code \u2191", FONT_BODY, Pt(14), SLATE_400, alignment=PP_ALIGN.LEFT)
    add_text_box(slide, map_left - Inches(0.1), map_top + map_h - Inches(0.05), map_w / 2, Inches(0.3),
                 "Static / Code-only \u2193", FONT_BODY, Pt(14), SLATE_400, alignment=PP_ALIGN.LEFT)

    # Competitor dots
    competitors = [
        ("OpenFisca", map_left + Inches(1.0), map_top + map_h / 2 + Inches(1.5), SLATE_400),
        ("EUROMOD", map_left + Inches(1.8), map_top + map_h / 2 + Inches(1.0), SLATE_400),
        ("PolicyEngine", map_left + Inches(1.5), map_top + map_h / 2 - Inches(0.5), SLATE_400),
        ("ResStock", map_left + map_w / 2 + Inches(1.0), map_top + map_h / 2 + Inches(1.2), SLATE_400),
        ("demod", map_left + map_w / 2 + Inches(0.5), map_top + map_h / 2 + Inches(1.8), SLATE_400),
        ("LIAM2 \u2020", map_left + Inches(2.5), map_top + map_h / 2 + Inches(0.3), SLATE_300),
    ]

    for name, cx, cy, color in competitors:
        dot_size = Inches(0.3)
        add_circle(slide, cx - dot_size / 2, cy - dot_size / 2, dot_size, color, color)
        add_text_box(slide, cx + Inches(0.2), cy - Inches(0.15), Inches(1.5), Inches(0.3),
                     name, FONT_BODY, Pt(14), SLATE_500)

    # ReformLab dot (prominent)
    rl_x = map_left + map_w / 2 + Inches(2.0)
    rl_y = map_top + Inches(1.0)
    rl_dot = Inches(0.5)
    add_circle(slide, rl_x - rl_dot / 2, rl_y - rl_dot / 2, rl_dot, BLUE_500, BLUE_500)
    add_text_box(slide, rl_x + Inches(0.35), rl_y - Inches(0.2), Inches(2.0), Inches(0.4),
                 "ReformLab", FONT_HEADING, Pt(18), BLUE_500, bold=True)

    # Legend on the right
    legend_x = Inches(9.2)
    legend_top = CONTENT_TOP + Inches(0.5)
    legend_items = [
        ("OpenFisca, EUROMOD", "Strong on income/tax rules"),
        ("PolicyEngine", "Public-facing tax-benefit delivery"),
        ("ResStock, demod", "Physical models only"),
        ("LIAM2, OpenM++", "Discontinued / archived"),
    ]

    add_text_box(slide, legend_x, legend_top - Inches(0.3), Inches(3.5), Inches(0.3),
                 "Key:", FONT_HEADING, Pt(16), SLATE_700, bold=True)

    for i, (name, desc) in enumerate(legend_items):
        ly = legend_top + Inches(0.1) + i * Inches(0.65)
        add_text_box(slide, legend_x, ly, Inches(3.5), Inches(0.3),
                     name, FONT_BODY, Pt(14), SLATE_700, bold=True)
        add_text_box(slide, legend_x, ly + Inches(0.25), Inches(3.5), Inches(0.3),
                     desc, FONT_BODY, Pt(14), SLATE_500)

    # Bottom conclusion
    bottom_y = map_top + map_h + Inches(0.6)
    add_text_box(
        slide, Inches(1.0), bottom_y, Inches(11), Inches(0.5),
        "Environmental + distributional + dynamic + no-code: no one else occupies this space.",
        FONT_HEADING, Pt(18), SLATE_900, bold=True, alignment=PP_ALIGN.CENTER
    )

    # Dagger footnote
    add_text_box(slide, Inches(1.0), bottom_y + Inches(0.5), Inches(5), Inches(0.3),
                 "\u2020 Discontinued", FONT_BODY, Pt(14), SLATE_400)

    add_speaker_notes(slide,
        "Key talking point: Walk through the positioning map. Tax-benefit tools are bottom-left (static, code-only, income-focused). Energy tools are bottom-right (environmental but no distribution). Legacy tools are being retired. ReformLab is alone in the top-right: environmental, distributional, dynamic, and no-code.\n\n"
        "Emotional beat: Confidence — 'No one else is doing this.'\n\n"
        "Transition: 'And the timing couldn't be better.'")

    # =========================================================================
    # SLIDE 12: WHY NOW
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_slide_title(slide, "Three forces converging right now")

    forces = [
        ("\U0001F3DB\uFE0F", "The Policy Moment",
         "EU Green Deal, carbon border adjustment,\nnational carbon taxes. Governments need\nassessment tools faster than ever.",
         BLUE_500),
        ("\U0001F6E0\uFE0F", "The Tool Gap",
         "LIAM2 discontinued. OpenM++ archived.\nNo replacement serves the environmental\npolicy workflow.",
         VIOLET_500),
        ("\u26A1", "The Technology Shift",
         "Mature open-source engines, powerful\nscientific Python stack, AI-assisted\ndevelopment. Buildable now.",
         EMERALD_500),
    ]

    col_w = Inches(3.5)
    col_h = Inches(4.2)
    col_gap = Inches(0.5)
    total_cols_w = 3 * col_w + 2 * col_gap
    cols_left = (SLIDE_WIDTH - total_cols_w) / 2
    cols_top = CONTENT_TOP + Inches(0.4)

    for i, (icon, title, desc, color) in enumerate(forces):
        x = cols_left + i * (col_w + col_gap)

        # Card
        add_rounded_rect(slide, x, cols_top, col_w, col_h, WHITE, border_color=SLATE_200)

        # Top color bar
        add_rect(slide, x, cols_top, col_w, Inches(0.06), color, border=False)

        # Icon
        add_text_box(slide, x, cols_top + Inches(0.4), col_w, Inches(0.7),
                     icon, FONT_HEADING, Pt(40), color, alignment=PP_ALIGN.CENTER)

        # Title
        add_text_box(slide, x + Inches(0.3), cols_top + Inches(1.3), col_w - Inches(0.6), Inches(0.5),
                     title, FONT_HEADING, Pt(20), SLATE_900, bold=True, alignment=PP_ALIGN.CENTER)

        # Description
        add_text_box(slide, x + Inches(0.3), cols_top + Inches(2.0), col_w - Inches(0.6), Inches(1.8),
                     desc, FONT_BODY, Pt(16), SLATE_500, alignment=PP_ALIGN.CENTER)

    add_speaker_notes(slide,
        "Key talking point: Three forces converging: urgent policy demand, dying legacy tools, and mature technology. This window won't stay open forever.\n\n"
        "Emotional beat: Urgency — 'The timing is right.'\n\n"
        "Transition: 'Ready to see it in action?'")

    # =========================================================================
    # SLIDE 13: GET STARTED / CLOSING
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    # Accent line at top (mirrors slide 1)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), BLUE_500)

    # Product name
    add_text_box(
        slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
        "ReformLab", FONT_HEADING, Pt(48), SLATE_900, bold=True,
        alignment=PP_ALIGN.CENTER
    )

    # Tagline
    add_text_box(
        slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.8),
        "See the impact before the vote.", FONT_HEADING, Pt(28), BLUE_500,
        alignment=PP_ALIGN.CENTER, italic=True
    )

    # Value props in a row
    props = [
        ("\U0001F513", "Open-source\nNo vendor lock-in"),
        ("\U0001F4BB", "Works offline\nNo cloud dependency"),
        ("\U0001F4CA", "Open-data default\nFunctional out of the box"),
        ("\u23F1\uFE0F", "30 minutes\nInstall to first output"),
    ]

    prop_w = Inches(2.5)
    prop_h = Inches(1.6)
    prop_gap = Inches(0.3)
    total_props_w = 4 * prop_w + 3 * prop_gap
    props_left = (SLIDE_WIDTH - total_props_w) / 2
    props_top = Inches(4.2)

    for i, (icon, text) in enumerate(props):
        x = props_left + i * (prop_w + prop_gap)
        add_rounded_rect(slide, x, props_top, prop_w, prop_h, SLATE_50, border_color=SLATE_200)
        add_text_box(slide, x, props_top + Inches(0.15), prop_w, Inches(0.5),
                     icon, FONT_HEADING, Pt(28), BLUE_500, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), props_top + Inches(0.65), prop_w - Inches(0.4), Inches(0.8),
                     text, FONT_BODY, Pt(16), SLATE_700, alignment=PP_ALIGN.CENTER)

    # CTA
    add_text_box(
        slide, Inches(3.0), Inches(6.2), Inches(7), Inches(0.5),
        "Let's talk \u2014 reformlab.org", FONT_HEADING, Pt(22), BLUE_500,
        bold=True, alignment=PP_ALIGN.CENTER
    )

    add_speaker_notes(slide,
        "Key talking point: Recap the four key value props — open source, offline, open data, fast to start. The audience should leave thinking 'I want to try this.'\n\n"
        "Emotional beat: Action — 'I want to try this.'\n\n"
        "Close with: 'Thank you. I'm happy to answer questions or show you a live demo.'")

    # Save
    output_path = os.path.join(os.path.dirname(__file__), "ReformLab-Pitch-Deck.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


# ==============================================================================
# HELPER FUNCTIONS
# ==============================================================================

def set_slide_bg(slide, color):
    """Set slide background to solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_name, font_size,
                 font_color, bold=False, italic=False, alignment=PP_ALIGN.LEFT,
                 vertical=MSO_ANCHOR.TOP):
    """Add a text box with specified formatting."""
    from pptx.util import Emu as EmuConv
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    # Set vertical alignment
    tf.paragraphs[0].alignment = alignment

    # Handle multi-line text
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = font_size
        run.font.color.rgb = font_color
        run.font.bold = bold
        run.font.italic = italic
        p.space_before = Pt(0)
        p.space_after = Pt(2)

    # Vertical anchoring
    txBox.text_frame.paragraphs[0].alignment = alignment
    from pptx.oxml.ns import qn
    bodyPr = txBox.text_frame._txBody.find(qn('a:bodyPr'))
    if vertical == MSO_ANCHOR.MIDDLE:
        bodyPr.set('anchor', 'ctr')
    elif vertical == MSO_ANCHOR.BOTTOM:
        bodyPr.set('anchor', 'b')

    return txBox


def add_rect(slide, left, top, width, height, fill_color, border=True, border_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border and border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Adjust corner radius
    shape.adjustments[0] = 0.02  # Slight rounding
    return shape


def add_circle(slide, left, top, size, fill_color, border_color=None):
    """Add a circle (oval) shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_slide_title(slide, text):
    """Add a consistent slide title."""
    add_text_box(
        slide, Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.8),
        text, FONT_HEADING, Pt(28), SLATE_900, bold=True
    )
    # Subtle underline
    add_rect(slide, Inches(0.6), Inches(1.15), Inches(2.0), Inches(0.04), BLUE_500, border=False)


def add_speaker_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text


if __name__ == "__main__":
    create_presentation()
